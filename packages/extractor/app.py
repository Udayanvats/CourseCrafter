import pika
from pptx import Presentation      
import boto3
import os
from dotenv import load_dotenv,dotenv_values
import tempfile 
import json
from PyPDF2 import PdfReader
import pytesseract
from PIL import Image
import io



load_dotenv()


AWS_SERVER_PUBLIC_KEY = os.getenv('AWS_ACCESS_KEY')
AWS_SERVER_SECRET_KEY = os.getenv('AWS_SECRET_KEY')
HOST = os.getenv('HOST') or 'localhost'
print(AWS_SERVER_PUBLIC_KEY)
print(AWS_SERVER_SECRET_KEY)

print(AWS_SERVER_PUBLIC_KEY, AWS_SERVER_SECRET_KEY)

s3 = boto3.client('s3',
    aws_access_key_id=AWS_SERVER_PUBLIC_KEY, 
    aws_secret_access_key=AWS_SERVER_SECRET_KEY, 
    region_name="ap-south-1"
)

RABBITMQ_HOST = HOST
RABBITMQ_PORT = 5672
RABBITMQ_USERNAME = os.getenv('RABBITMQ_USER')
RABBITMQ_PASSWORD = os.getenv('RABBITMQ_PASSWORD')
bucket_name = 'coursecrafter'

credentials = pika.PlainCredentials(RABBITMQ_USERNAME, RABBITMQ_PASSWORD)
parameters = pika.ConnectionParameters(host=RABBITMQ_HOST, port=RABBITMQ_PORT, credentials=credentials)

connection = pika.BlockingConnection(parameters)
channel = connection.channel()


channel.queue_declare(queue='extract')
channel.queue_declare(queue='notification',durable=True)

def notify_user(body):
    print("Notifying user")

    channel.basic_publish(exchange='', routing_key='notification', body=body)

def upload_file_to_s3(file,object_name):
   
    s3.upload_fileobj(file,bucket_name,object_name)

def get_file_from_s3(filepath):
    try:
     
        print("Downloading ", filepath)
        filename = filepath.split("/")[-1]
        
        directory = "uploads/"
        
        if not os.path.exists(directory):
            os.makedirs(directory)
        
        full_path = os.path.join(directory, filename)

        print(f"filepath  {filepath}")
        
        with open(full_path, "wb") as f:
            s3.download_fileobj(bucket_name, filepath, f)

        print(f"File downloaded to {full_path}")

        return full_path
    except Exception as e:
        print(f"Error downloading file from S3: {e}")
        return None



def extract_text_from_ppt(ppt_path):
    ppt = Presentation(ppt_path)
    text =""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text+=shape.text
                text+="\n"

    
    return text

def extract_text_from_pdf(filepath):
    reader = PdfReader(filepath)
    number_of_pages = len(reader.pages)
    text=""
    for i in range(number_of_pages):
        page = reader.pages[i]
        text += page.extract_text()
        for image in page.images:
            # file=os.write(tempfile.NamedTemporaryFile(suffix=".png",delete=False),image.data)

            text+=extract_text_from_image(Image.open(io.BytesIO(image.data)))
    
    print(text,"text")
    return text


def extract_text_from_image(image):
    text = pytesseract.image_to_string(image)
    return text



def extract_text(filepath):
    if filepath.endswith(".pptx"):
        return extract_text_from_ppt(filepath)
    elif filepath.endswith(".pdf"):
        return extract_text_from_pdf(filepath)
    else:
        return None
    



def process_documents(ch, method, properties, body):
    try:
        reqObject = body.decode('utf-8')
        print(reqObject)
        json_object = json.loads(reqObject)
        docs = json_object["docs"]
        courseId = json_object["courseId"]
        mode= json_object["mode"]
        print("GOt course id", courseId)
        pyqs = json_object["pyqs"]
        print(f"docs: {docs}")
        print(f"pyqs: {pyqs}")
        print(f"courseId: {courseId}")
        json_data = {
            "pyqs": [],
            "docs": []
        }


        if docs:
            for doc in docs:
                s3ObjPath = doc
                filepath = get_file_from_s3(s3ObjPath)
                fileid = s3ObjPath.split("/")[-1].split(".")[0]
                filename = s3ObjPath.split("/")[-1]
                text = extract_text(filepath)
                os.remove(filepath)
                json_data["docs"].append({
                    "filename": filename,
                    "contents": text
                })
                notify_user(json.dumps({
                    "status": True,
                    "error": "",
                    "courseId": courseId,
                    "message": filename,
                }))

          

        if pyqs:
            for pyq in pyqs:
                s3ObjPath = pyq
                filepath = get_file_from_s3(s3ObjPath)
                print(f"file path: {filepath}")
                fileid = s3ObjPath.split("/")[-1].split(".")[0]
                text = extract_text(filepath)
                os.remove(filepath)

                notify_user(json.dumps({
                    "status": True,
                    "error": "",
                    "courseId": courseId,
                    "message": pyq,
    
                }))
                json_data["pyqs"].append({
                    "filename": pyq,
                    "contents": text
                })

        

        with tempfile.NamedTemporaryFile(mode="w", encoding="utf-8", delete=False) as temp_file:
            temp_file.write(json.dumps(json_data))
            temp_file_name = temp_file.name
        print(f"Temp file created: {temp_file_name}")
        upload_file_to_s3(open(temp_file_name, "rb"), "text/" + courseId + ".json")
        os.unlink(temp_file_name)  

        notify_user(json.dumps({
            "status": True,
            "object_path": "text/" + courseId + ".json",
            "error": "",
            "message": "[DONE]",
            "courseId": courseId,
            "mode": mode
        }))

    except Exception as e:
        notify_user(json.dumps({
            "status": False,
            "object_path": "",
            "error": "Error extracting text: " + str(e)
        }))
        print(f"Error extracting text: {e}")


channel.basic_consume(queue='extract', on_message_callback=process_documents, auto_ack=True)

print('Waiting for messages')

if __name__ == '__main__':
    print("Starting extractor")
    channel.start_consuming()
    



